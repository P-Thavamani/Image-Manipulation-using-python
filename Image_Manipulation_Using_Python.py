import tkinter as tk
from tkinter import filedialog
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from docx.shared import Inches
import os

# Conversion Functions

def convert_to_pdf():
    image_path = file_path_entry.get()
    if image_path:
        with Image.open(image_path) as img:
            pdf_path = "converted_to_pdf.pdf"
            img.save(os.path.join(output_dir, pdf_path))
            status_label.config(text="Image converted to PDF: {}".format(os.path.abspath(os.path.join(output_dir, pdf_path))))
    else:
        status_label.config(text="Please select an image file.")

def convert_to_ppt():
    image_path = file_path_entry.get()
    if image_path:
        with Image.open(image_path) as img:
            ppt_path = "converted_to_ppt.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide_width, slide_height = prs.slide_width, prs.slide_height
            max_width, max_height = slide_width - (2 * Inches(0.5)), slide_height - (2 * Inches(0.5))
            img.thumbnail((max_width, max_height))
            img_width, img_height = img.size
            img_width *= 5735
            img_height *= 5735
            left, top = (slide_width - img_width) / 2, (slide_height - img_height) / 2
            slide.shapes.add_picture(image_path, left, top, width=img_width, height=img_height)
            prs.save(os.path.join(output_dir, ppt_path))
            status_label.config(text="Image converted to PPT: {}".format(os.path.abspath(os.path.join(output_dir, ppt_path))))
    else:
        status_label.config(text="Please select an image file.")

def convert_to_jpg():
    image_path = file_path_entry.get()
    if image_path:
        with Image.open(image_path) as img:
            jpg_path = "converted_to_jpg.jpg"
            if img.mode == "RGBA":
                img = img.convert("RGB")
            img.save(os.path.join(output_dir, jpg_path))
            status_label.config(text="Image converted to JPG: {}".format(os.path.abspath(os.path.join(output_dir, jpg_path))))
    else:
        status_label.config(text="Please select an image file.")

def convert_to_png():
    image_path = file_path_entry.get()
    if image_path:
        with Image.open(image_path) as img:
            png_path = "converted_to_png.png"
            img.save(os.path.join(output_dir, png_path))
            status_label.config(text="Image converted to PNG: {}".format(os.path.abspath(os.path.join(output_dir, png_path))))
    else:
        status_label.config(text="Please select an image file.")

def convert_to_word():
    image_path = file_path_entry.get()
    if image_path:
        with Image.open(image_path) as img:
            doc_path = "converted_to_word.docx"
            document = Document()
            document.add_picture(image_path, width=Inches(6))
            document.save(os.path.join(output_dir, doc_path))
            status_label.config(text="Image converted to Word: {}".format(os.path.abspath(os.path.join(output_dir, doc_path))))
    else:
        status_label.config(text="Please select an image file.")

# Image Manipulation

def resize_image():
    image_path = file_path_entry.get()
    if image_path:
        try:
            width = int(width_entry.get())
            height = int(height_entry.get())
            with Image.open(image_path) as img:
                resized_img = img.resize((width, height))
                resized_img = resized_img.convert("RGB")  # Convert to RGB color mode
                resized_path = "resized_image.jpg"
                resized_img.save(os.path.join(output_dir, resized_path))
                status_label.config(text="Image resized: {}".format(os.path.abspath(os.path.join(output_dir, resized_path))))
        except ValueError:
            status_label.config(text="Please enter valid width and height values.")
    else:
        status_label.config(text="Please select an image file.")

# File Selection

def select_file():
    file_path = filedialog.askopenfilename(filetypes=[("Image Files", "*.jpg;*.jpeg;*.png")])
    file_path_entry.delete(0, tk.END)
    file_path_entry.insert(tk.END, file_path)

def select_output_directory():
    global output_dir
    output_dir = filedialog.askdirectory()
    output_dir_label.config(text="Output Directory: " + output_dir)

# Main Application

root = tk.Tk()
root.title("Image Resizer and File Converter")
root.geometry("700x600")
root.resizable(False, False)

# Widgets

file_path_label = tk.Label(root, text="Image File:", bg="white", font=("Arial", 12, "bold"))
file_path_label.pack()

file_path_entry = tk.Entry(root, width=40, font=("Arial", 12))
file_path_entry.pack()

select_file_button = tk.Button(root, text="Select File", command=select_file, font=("Arial", 12))
select_file_button.pack()

output_dir_label = tk.Label(root, text="Output Directory: ", bg="white", font=("Arial", 12, "bold"))
output_dir_label.pack()

select_output_dir_button = tk.Button(root, text="Select Output Directory", command=select_output_directory, font=("Arial", 12))
select_output_dir_button.pack()

width_label = tk.Label(root, text="Width:", bg="white", font=("Arial", 12, "bold"))
width_label.pack()

width_entry = tk.Entry(root, width=10, font=("Arial", 12))
width_entry.pack()

height_label = tk.Label(root, text="Height:", bg="white", font=("Arial", 12, "bold"))
height_label.pack()

height_entry = tk.Entry(root, width=10, font=("Arial", 12))
height_entry.pack()

resize_button = tk.Button(root, text="Resize Image", command=resize_image, font=("Arial", 12))
resize_button.pack(pady=10)

convert_to_pdf_button = tk.Button(root, text="Convert to PDF", command=convert_to_pdf, font=("Arial", 12))
convert_to_pdf_button.pack()

convert_to_ppt_button = tk.Button(root, text="Convert to PPT", command=convert_to_ppt, font=("Arial", 12))
convert_to_ppt_button.pack()

convert_to_jpg_button = tk.Button(root, text="Convert to JPG", command=convert_to_jpg, font=("Arial", 12))
convert_to_jpg_button.pack()

convert_to_png_button = tk.Button(root, text="Convert to PNG", command=convert_to_png, font=("Arial", 12))
convert_to_png_button.pack()

convert_to_word_button = tk.Button(root, text="Convert to Word", command=convert_to_word, font=("Arial", 12))
convert_to_word_button.pack()

status_label = tk.Label(root, text="", bg="white", font=("Arial", 12))
status_label.pack()

root.mainloop()
